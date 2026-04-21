"""Document to markdown converter.

Converts PDF, DOCX, PPTX, XLSX, XLS, HTML, CSV, TXT, RTF, MD to markdown
with YAML frontmatter ready for vault ingestion.

Design:
  - Script-first, zero LLM calls. Pure extraction + formatting.
  - SHA256-keyed dedup via IngestLedger.
  - Source file is NEVER copied into the vault. The frontmatter
    `source_file` field points back to the original location.

Ported and generalized from Sean Uddin's MKV scripts.
"""

from __future__ import annotations

import hashlib
import os
import re
from dataclasses import dataclass, field
from datetime import datetime
from typing import Callable

from mymory.core.ledger import IngestLedger


# ======================================================================
# Format handlers
# ======================================================================


def convert_pdf(filepath: str) -> str:
    """Convert PDF to markdown using PyMuPDF."""
    import fitz  # pymupdf

    doc = fitz.open(filepath)
    pages: list[str] = []
    has_text = False

    for i, page in enumerate(doc):
        text = page.get_text("text")
        if text.strip():
            has_text = True
            pages.append(f"<!-- Page {i+1} -->\n\n{text.strip()}")

    doc.close()

    if not has_text:
        return (
            "> **OCR REQUIRED**: This PDF contains no extractable text.\n"
            "> Route through an OCR pipeline (Marker + Surya, Tesseract, etc.).\n"
        )

    content = "\n\n---\n\n".join(pages)
    content = re.sub(r"\n{4,}", "\n\n\n", content)
    content = re.sub(r" {3,}", "  ", content)
    # Fix hyphenated line breaks: "word-\nbreak" -> "wordbreak"
    content = re.sub(r"(\w)-\n(\w)", r"\1\2", content)
    return content


def convert_docx(filepath: str) -> str:
    """Convert DOCX to markdown. Primary via mammoth, fallback python-docx."""
    md_content = ""
    try:
        import mammoth

        with open(filepath, "rb") as f:
            result = mammoth.convert_to_markdown(f)
            md_content = result.value
    except ImportError:
        pass

    if len(md_content.strip()) < 50:
        import docx

        doc = docx.Document(filepath)
        paragraphs: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            style = para.style.name.lower() if para.style else ""
            if "heading 1" in style or "title" in style:
                paragraphs.append(f"# {text}")
            elif "heading 2" in style:
                paragraphs.append(f"## {text}")
            elif "heading 3" in style:
                paragraphs.append(f"### {text}")
            else:
                paragraphs.append(text)

        for table in doc.tables:
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
                rows.insert(1, header_sep)
                paragraphs.append("\n".join(rows))

        md_content = "\n\n".join(paragraphs)

    return re.sub(r"\n{4,}", "\n\n\n", md_content)


def convert_pptx(filepath: str) -> str:
    """Convert PPTX to markdown: one section per slide, speaker notes included."""
    from pptx import Presentation

    prs = Presentation(filepath)
    slides: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"## Slide {i}"]

        title_text = ""
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip()
            parts[0] = f"## Slide {i}: {title_text}"

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and text != title_text:
                    parts.append(text)
            if shape.has_table:
                table = shape.table
                rows: list[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    header_sep = "| " + " | ".join(["---"] * len(table.columns)) + " |"
                    rows.insert(1, header_sep)
                    parts.append("\n".join(rows))

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"\n> **Speaker Notes:** {notes}")

        slides.append("\n\n".join(parts))

    return "\n\n---\n\n".join(slides)


def convert_xlsx(filepath: str) -> str:
    """Convert XLSX to markdown tables, one section per sheet."""
    import openpyxl

    wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
    sheets: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if not any(cells):
                continue
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            ncols = rows[0].count("|") - 1
            rows.insert(1, "| " + " | ".join(["---"] * ncols) + " |")
            sheets.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(rows))

    wb.close()
    return "\n\n---\n\n".join(sheets) if sheets else "> Empty spreadsheet"


def convert_xls(filepath: str) -> str:
    """Convert legacy .xls via pandas + xlrd."""
    import pandas as pd

    try:
        sheets = pd.read_excel(filepath, sheet_name=None, engine="xlrd")
    except Exception:
        try:
            sheets = pd.read_excel(filepath, sheet_name=None)
        except Exception as e:
            return f"> **CONVERSION FAILED**: {e}"

    parts: list[str] = []
    for name, df in sheets.items():
        parts.append(f"## Sheet: {name}\n\n{df.to_markdown(index=False)}")
    return "\n\n---\n\n".join(parts) if parts else "> Empty spreadsheet"


def convert_html(filepath: str) -> str:
    """Convert HTML to markdown via BeautifulSoup with structure preservation."""
    from bs4 import BeautifulSoup

    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()

    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    parts: list[str] = []
    title_tag = soup.find("title")
    if title_tag:
        title = title_tag.get_text().strip()
        if title:
            parts.append(f"# {title}")

    for element in soup.find_all(
        ["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "pre", "code", "table", "blockquote"]
    ):
        tag = element.name
        text = element.get_text().strip()
        if not text:
            continue
        if tag == "h1":
            parts.append(f"# {text}")
        elif tag == "h2":
            parts.append(f"## {text}")
        elif tag == "h3":
            parts.append(f"### {text}")
        elif tag in ("h4", "h5", "h6"):
            parts.append(f"#### {text}")
        elif tag in ("pre", "code"):
            parts.append(f"```\n{text}\n```")
        elif tag == "li":
            parts.append(f"- {text}")
        elif tag == "blockquote":
            parts.append(f"> {text}")
        elif tag == "table":
            rows = element.find_all("tr")
            if rows:
                md_rows: list[str] = []
                for row in rows:
                    cells = [c.get_text().strip() for c in row.find_all(["td", "th"])]
                    md_rows.append("| " + " | ".join(cells) + " |")
                if len(md_rows) > 1:
                    ncols = md_rows[0].count("|") - 1
                    md_rows.insert(1, "| " + " | ".join(["---"] * ncols) + " |")
                parts.append("\n".join(md_rows))
        elif tag == "p":
            parts.append(text)

    content = "\n\n".join(parts)

    # Dedup consecutive identical paragraphs (common in HTML extraction).
    lines = content.split("\n\n")
    deduped: list[str] = []
    prev: str | None = None
    for line in lines:
        if line != prev:
            deduped.append(line)
        prev = line
    content = "\n\n".join(deduped)
    return re.sub(r"\n{4,}", "\n\n\n", content)


def convert_csv(filepath: str) -> str:
    """Convert CSV to markdown table via pandas."""
    try:
        import pandas as pd

        df = pd.read_csv(filepath)
        return df.to_markdown(index=False)
    except Exception:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()[:200]
        return "```csv\n" + "".join(lines) + "\n```"


def convert_txt(filepath: str) -> str:
    """TXT passthrough. RTF stripped of basic control codes."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()

    if filepath.lower().endswith(".rtf"):
        try:
            from striprtf.striprtf import rtf_to_text

            content = rtf_to_text(content)
        except ImportError:
            # Lightweight fallback: strip RTF control sequences.
            content = re.sub(r"\\[a-z]+\d*\s?", "", content)
            content = re.sub(r"[{}]", "", content)
            content = content.strip()

    return content


def convert_md(filepath: str) -> str:
    """Markdown passthrough (strips any existing frontmatter to avoid conflicts)."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    m = re.match(r"^---\s*\n.*?\n---\s*\n?(.*)$", content, re.DOTALL)
    return m.group(1) if m else content


# ======================================================================
# Dispatch
# ======================================================================


CONVERTERS: dict[str, Callable[[str], str]] = {
    ".pdf": convert_pdf,
    ".docx": convert_docx,
    ".pptx": convert_pptx,
    ".xlsx": convert_xlsx,
    ".xls": convert_xls,
    ".html": convert_html,
    ".htm": convert_html,
    ".csv": convert_csv,
    ".txt": convert_txt,
    ".rtf": convert_txt,
    ".md": convert_md,
    ".markdown": convert_md,
}


# ======================================================================
# Helpers
# ======================================================================


def file_sha256(filepath: str) -> str:
    h = hashlib.sha256()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def slugify(name: str, max_len: int = 80) -> str:
    name = os.path.splitext(name)[0]
    slug = re.sub(r"[^\w\-]", "_", name)
    slug = re.sub(r"_+", "_", slug).strip("_").lower()
    return slug[:max_len]


def extract_title(content: str, filename: str) -> str:
    m = re.search(r"^#\s+(.+)$", content, re.MULTILINE)
    if m:
        title = m.group(1).strip()
    else:
        title = os.path.splitext(filename)[0].replace("_", " ").replace("-", " ")
        title = " ".join(w if w.isupper() else w.title() for w in title.split())
    return title.replace('"', "'")[:120]


# ======================================================================
# Conversion result + API
# ======================================================================


@dataclass
class ConvertResult:
    success: bool = False
    source: str = ""
    dest_path: str | None = None
    word_count: int = 0
    skipped: bool = False
    error: str | None = None
    sha256: str | None = None

    def to_dict(self) -> dict:
        return {
            "success": self.success,
            "source": self.source,
            "dest_path": self.dest_path,
            "word_count": self.word_count,
            "skipped": self.skipped,
            "error": self.error,
            "sha256": self.sha256,
        }


def build_frontmatter(
    title: str,
    wing: str,
    created: str,
    source_file: str,
    source_format: str,
    source_size: int,
    sha256: str,
    word_count: int,
    tags: list[str],
    confidence: str = "VERBATIM",
    source_agent: str = "mymory-converter",
) -> str:
    """Generate YAML frontmatter block (no wrapping --- pairs: returned with them)."""
    tags_block = "\n".join(f"  - {t}" for t in tags) if tags else "  - imported"
    return (
        "---\n"
        f'title: "{_escape(title)}"\n'
        f"wing: {wing}\n"
        f"created: {created}\n"
        f"imported: {datetime.now().strftime('%Y-%m-%d')}\n"
        f'source_file: "{_escape(source_file)}"\n'
        f"source_format: {source_format}\n"
        f"source_size: {source_size}\n"
        f"source_sha256: {sha256}\n"
        f"source_agent: {source_agent}\n"
        f"confidence: {confidence}\n"
        f"word_count: {word_count}\n"
        "tags:\n"
        f"{tags_block}\n"
        "---\n\n"
    )


def _escape(s: str) -> str:
    return s.replace("\\", "\\\\").replace('"', '\\"')


def convert_file(
    filepath: str,
    vault_root: str,
    wing: str,
    ledger: IngestLedger | None = None,
    tags: list[str] | None = None,
    out_path: str | None = None,
    confidence: str = "VERBATIM",
) -> ConvertResult:
    """Convert a single file. Writes into `vault_root/wing/` unless `out_path` is given.

    Returns a ConvertResult capturing outcome, destination path, and stats.
    """
    r = ConvertResult(source=filepath)

    if not os.path.isfile(filepath):
        r.error = f"file not found: {filepath}"
        return r

    ext = os.path.splitext(filepath)[1].lower()
    if ext not in CONVERTERS:
        r.error = f"unsupported format: {ext}"
        return r

    sha = file_sha256(filepath)
    r.sha256 = sha

    if ledger and ledger.is_converted(sha):
        r.skipped = True
        r.error = "already converted (sha256 match)"
        return r

    stat = os.stat(filepath)
    mtime = datetime.fromtimestamp(stat.st_mtime)
    date_str = mtime.strftime("%Y-%m-%d")
    fname = os.path.basename(filepath)

    try:
        content = CONVERTERS[ext](filepath)
    except Exception as e:
        r.error = f"conversion failed: {type(e).__name__}: {e}"
        return r

    if not content or len(content.strip()) < 10:
        r.error = "empty or near-empty output"
        return r

    title = extract_title(content, fname)
    word_count = len(content.split())

    fm = build_frontmatter(
        title=title,
        wing=wing,
        created=date_str,
        source_file=filepath,
        source_format=ext.lstrip("."),
        source_size=stat.st_size,
        sha256=sha,
        word_count=word_count,
        tags=tags or ["imported"],
        confidence=confidence,
    )

    if out_path:
        dest_path = out_path
        os.makedirs(os.path.dirname(os.path.abspath(dest_path)), exist_ok=True)
    else:
        slug = slugify(fname)
        wing_dir = os.path.join(vault_root, wing)
        os.makedirs(wing_dir, exist_ok=True)
        dest_path = os.path.join(wing_dir, f"{date_str}_{slug}.md")
        counter = 1
        while os.path.exists(dest_path):
            dest_path = os.path.join(wing_dir, f"{date_str}_{slug}_{counter}.md")
            counter += 1

    # Sanitise inline `#` tokens so hex colors, addresses, and content
    # hashtags in source documents don't pollute the Obsidian tag graph.
    from mymory.core.sanitise import sanitise_hashes
    content = sanitise_hashes(content)

    with open(dest_path, "w", encoding="utf-8") as f:
        f.write(fm)
        f.write(content)

    if ledger:
        ledger.record(
            sha256=sha,
            source_path=filepath,
            dest_path=dest_path,
            wing=wing,
            source_format=ext.lstrip("."),
            word_count=word_count,
        )

    r.success = True
    r.dest_path = dest_path
    r.word_count = word_count
    return r


def convert_batch(
    files: list[str],
    vault_root: str,
    wing: str,
    ledger: IngestLedger | None = None,
    tags: list[str] | None = None,
) -> dict:
    """Convert a batch of files. Returns summary."""
    stats = {
        "converted": 0,
        "skipped": 0,
        "failed": 0,
        "total_words": 0,
        "errors": [],
        "results": [],
    }

    for fp in files:
        res = convert_file(fp, vault_root, wing, ledger=ledger, tags=tags)
        stats["results"].append(res.to_dict())
        if res.success:
            stats["converted"] += 1
            stats["total_words"] += res.word_count
        elif res.skipped:
            stats["skipped"] += 1
        else:
            stats["failed"] += 1
            stats["errors"].append({"file": fp, "error": res.error})

    return stats


def supported_extensions() -> list[str]:
    return sorted(CONVERTERS.keys())
